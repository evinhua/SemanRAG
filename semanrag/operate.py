"""SemanRAG – Document processing and query pipeline."""

from __future__ import annotations

import asyncio
import json
import re
from datetime import datetime
from typing import Any

import numpy as np

try:
    import json_repair
except ImportError:
    json_repair = None  # type: ignore[assignment]

from semanrag.base import (
    ACLPolicy,
    BaseGraphStorage,
    BaseKVStorage,
    BaseLexicalStorage,
    BaseVectorStorage,
    DocStatus,
    DocStatusStorage,
    ExtractionResult,
    ExtractedEntity,
    ExtractedRelation,
    QueryParam,
    QueryResult,
    TextChunkSchema,
)
from semanrag.prompt import (
    DEFAULT_COMPLETION_DELIMITER,
    DEFAULT_RECORD_DELIMITER,
    DEFAULT_TUPLE_DELIMITER,
    PROMPTS,
)
from semanrag.utils import (
    EmbeddingFunc,
    TiktokenTokenizer,
    compute_mdhash_id,
    detect_pii,
    detect_prompt_injection,
    logger,
    otel_span,
    reciprocal_rank_fusion,
    sanitize_output,
    truncate_list_by_token_size,
    use_llm_func_with_cache,
)


# ═══════════════════════════════════════════════════════════════════════════
# CHUNKING
# ═══════════════════════════════════════════════════════════════════════════


def chunking_by_token_size(
    content: str,
    overlap_token_size: int = 100,
    max_token_size: int = 1200,
    tokenizer=None,
    split_by_character: str | None = None,
    split_by_character_only: bool = False,
) -> list[TextChunkSchema]:
    if tokenizer is None:
        tokenizer = TiktokenTokenizer()
    results: list[TextChunkSchema] = []

    if split_by_character and split_by_character_only:
        segments = content.split(split_by_character)
        for idx, seg in enumerate(segments):
            seg = seg.strip()
            if not seg:
                continue
            tokens = tokenizer.encode(seg)
            results.append(
                TextChunkSchema(
                    tokens=len(tokens),
                    content=seg,
                    chunk_order_index=idx,
                    full_doc_id="",
                    section_path=None,
                    page_number=None,
                    modality="text",
                )
            )
        return results

    if split_by_character:
        segments = content.split(split_by_character)
        all_tokens: list[int] = []
        seg_boundaries: list[int] = []
        for seg in segments:
            toks = tokenizer.encode(seg)
            all_tokens.extend(toks)
            seg_boundaries.append(len(all_tokens))
    else:
        all_tokens = tokenizer.encode(content)

    if not all_tokens:
        return results

    idx = 0
    chunk_index = 0
    while idx < len(all_tokens):
        end = min(idx + max_token_size, len(all_tokens))
        chunk_tokens = all_tokens[idx:end]
        chunk_text = tokenizer.decode(chunk_tokens)
        results.append(
            TextChunkSchema(
                tokens=len(chunk_tokens),
                content=chunk_text.strip(),
                chunk_order_index=chunk_index,
                full_doc_id="",
                section_path=None,
                page_number=None,
                modality="text",
            )
        )
        chunk_index += 1
        if end >= len(all_tokens):
            break
        idx = end - overlap_token_size
        if idx <= (end - max_token_size):
            idx = end

    return results


def _split_sentences(text: str) -> list[str]:
    parts = re.split(r'(?<=[.!?])\s+', text)
    return [p.strip() for p in parts if p.strip()]


def chunking_semantic(
    content: str,
    embedding_func,
    drift_threshold: float = 0.5,
    min_size: int = 100,
    max_size: int = 2000,
) -> list[TextChunkSchema]:
    sentences = _split_sentences(content)
    if not sentences:
        return []

    loop = asyncio.get_event_loop()
    if loop.is_running():
        import concurrent.futures
        with concurrent.futures.ThreadPoolExecutor() as pool:
            embeddings = pool.submit(
                asyncio.run, embedding_func(sentences)
            ).result()
    else:
        embeddings = asyncio.run(embedding_func(sentences))

    embeddings = np.array(embeddings)
    split_points = [0]
    for i in range(1, len(embeddings)):
        a = embeddings[i - 1]
        b = embeddings[i]
        cos_sim = float(np.dot(a, b) / (np.linalg.norm(a) * np.linalg.norm(b) + 1e-10))
        if (1.0 - cos_sim) > drift_threshold:
            split_points.append(i)
    split_points.append(len(sentences))

    raw_chunks: list[str] = []
    for i in range(len(split_points) - 1):
        chunk_text = " ".join(sentences[split_points[i]:split_points[i + 1]])
        raw_chunks.append(chunk_text)

    merged: list[str] = []
    buf = ""
    for c in raw_chunks:
        if buf and len(buf) + len(c) > max_size:
            merged.append(buf)
            buf = c
        else:
            buf = (buf + " " + c).strip() if buf else c
    if buf:
        merged.append(buf)

    final: list[str] = []
    for c in merged:
        if len(c) > max_size:
            for j in range(0, len(c), max_size):
                part = c[j:j + max_size].strip()
                if part:
                    final.append(part)
        elif len(c) < min_size and final:
            final[-1] = final[-1] + " " + c
        else:
            final.append(c)

    tokenizer = TiktokenTokenizer()
    results: list[TextChunkSchema] = []
    for idx, chunk_text in enumerate(final):
        tokens = tokenizer.encode(chunk_text)
        results.append(
            TextChunkSchema(
                tokens=len(tokens),
                content=chunk_text,
                chunk_order_index=idx,
                full_doc_id="",
                section_path=None,
                page_number=None,
                modality="text",
            )
        )
    return results


def chunking_structure_aware(
    content: str,
    modality: str = "text",
    section_headers: list[str] | None = None,
) -> list[TextChunkSchema]:
    tokenizer = TiktokenTokenizer()
    results: list[TextChunkSchema] = []

    if modality == "text" and re.search(r'^#{1,6}\s', content, re.MULTILINE):
        header_pattern = re.compile(r'^(#{1,6})\s+(.+)$', re.MULTILINE)
        parts = header_pattern.split(content)
        sections: list[tuple[str, str]] = []
        path_stack: list[tuple[int, str]] = []
        i = 0
        preamble = parts[0].strip() if parts else ""
        if preamble:
            sections.append(("", preamble))
        i = 1
        while i < len(parts) - 2:
            level = len(parts[i])
            title = parts[i + 1].strip()
            body = parts[i + 2].strip() if i + 2 < len(parts) else ""
            while path_stack and path_stack[-1][0] >= level:
                path_stack.pop()
            path_stack.append((level, title))
            section_path = " > ".join(h for _, h in path_stack)
            if body:
                sections.append((section_path, body))
            i += 3

        for idx, (sec_path, sec_content) in enumerate(sections):
            tokens = tokenizer.encode(sec_content)
            results.append(
                TextChunkSchema(
                    tokens=len(tokens),
                    content=sec_content,
                    chunk_order_index=idx,
                    full_doc_id="",
                    section_path=sec_path or None,
                    page_number=None,
                    modality="text",
                )
            )
    else:
        paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
        for idx, para in enumerate(paragraphs):
            tokens = tokenizer.encode(para)
            results.append(
                TextChunkSchema(
                    tokens=len(tokens),
                    content=para,
                    chunk_order_index=idx,
                    full_doc_id="",
                    section_path=None,
                    page_number=None,
                    modality="text",
                )
            )

    return results


# ═══════════════════════════════════════════════════════════════════════════
# MULTI-MODAL PARSING
# ═══════════════════════════════════════════════════════════════════════════


async def parse_pdf(
    path: str,
    extract_tables: bool = True,
    extract_figures: bool = True,
    global_config: dict | None = None,
) -> list[TextChunkSchema]:
    results: list[TextChunkSchema] = []
    tokenizer = TiktokenTokenizer()
    global_config = global_config or {}

    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    tokens = tokenizer.encode(text)
                    results.append(
                        TextChunkSchema(
                            tokens=len(tokens),
                            content=text.strip(),
                            chunk_order_index=len(results),
                            full_doc_id="",
                            section_path=None,
                            page_number=page_num,
                            modality="text",
                        )
                    )
    except ImportError:
        try:
            from pypdf import PdfReader
            reader = PdfReader(path)
            for page_num, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    tokens = tokenizer.encode(text)
                    results.append(
                        TextChunkSchema(
                            tokens=len(tokens),
                            content=text.strip(),
                            chunk_order_index=len(results),
                            full_doc_id="",
                            section_path=None,
                            page_number=page_num,
                            modality="text",
                        )
                    )
        except ImportError:
            logger.warning("No PDF library available (pdfplumber or pypdf)")
            return results

    if extract_tables:
        try:
            import camelot
            tables = camelot.read_pdf(path, pages="all", flavor="lattice")
            for t_idx, table in enumerate(tables):
                df = table.df
                md_lines = ["| " + " | ".join(str(c) for c in df.columns) + " |"]
                md_lines.append("| " + " | ".join("---" for _ in df.columns) + " |")
                for _, row in df.iterrows():
                    md_lines.append("| " + " | ".join(str(v) for v in row) + " |")
                md_text = "\n".join(md_lines)
                tokens = tokenizer.encode(md_text)
                results.append(
                    TextChunkSchema(
                        tokens=len(tokens),
                        content=md_text,
                        chunk_order_index=len(results),
                        full_doc_id="",
                        section_path=f"table-{t_idx}",
                        page_number=table.page if hasattr(table, "page") else None,
                        modality="table",
                    )
                )
        except (ImportError, Exception) as e:
            logger.debug(f"Table extraction skipped: {e}")

    if extract_figures and global_config.get("vision_model_func"):
        try:
            import fitz  # PyMuPDF
            vision_func = global_config["vision_model_func"]
            doc = fitz.open(path)
            fig_idx = 0
            for page_num, page in enumerate(doc, start=1):
                images = page.get_images(full=True)
                for img_info in images:
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    if base_image:
                        import base64
                        img_bytes = base_image["image"]
                        img_b64 = base64.b64encode(img_bytes).decode()
                        surrounding = page.get_text()[:500]
                        caption_prompt = PROMPTS["figure_caption"].format(
                            surrounding_text=surrounding
                        )
                        caption = await vision_func(caption_prompt, image_data=img_b64)
                        tokens = tokenizer.encode(caption)
                        results.append(
                            TextChunkSchema(
                                tokens=len(tokens),
                                content=caption,
                                chunk_order_index=len(results),
                                full_doc_id="",
                                section_path=f"figure-{fig_idx}",
                                page_number=page_num,
                                modality="figure_caption",
                            )
                        )
                        fig_idx += 1
            doc.close()
        except (ImportError, Exception) as e:
            logger.debug(f"Figure extraction skipped: {e}")

    return results


async def parse_docx(path: str) -> list[TextChunkSchema]:
    results: list[TextChunkSchema] = []
    tokenizer = TiktokenTokenizer()
    try:
        from docx import Document
        doc = Document(path)
        for idx, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if not text:
                continue
            tokens = tokenizer.encode(text)
            results.append(
                TextChunkSchema(
                    tokens=len(tokens),
                    content=text,
                    chunk_order_index=idx,
                    full_doc_id="",
                    section_path=None,
                    page_number=None,
                    modality="text",
                )
            )
        for t_idx, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                rows.append([cell.text for cell in row.cells])
            if rows:
                md_lines = ["| " + " | ".join(rows[0]) + " |"]
                md_lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")
                for row in rows[1:]:
                    md_lines.append("| " + " | ".join(row) + " |")
                md_text = "\n".join(md_lines)
                tokens = tokenizer.encode(md_text)
                results.append(
                    TextChunkSchema(
                        tokens=len(tokens),
                        content=md_text,
                        chunk_order_index=len(results),
                        full_doc_id="",
                        section_path=f"table-{t_idx}",
                        page_number=None,
                        modality="table",
                    )
                )
    except ImportError:
        logger.warning("python-docx not available")
    return results


async def parse_pptx(path: str) -> list[TextChunkSchema]:
    results: list[TextChunkSchema] = []
    tokenizer = TiktokenTokenizer()
    try:
        from pptx import Presentation
        prs = Presentation(path)
        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        rows.append([cell.text for cell in row.cells])
                    if rows:
                        md_lines = ["| " + " | ".join(rows[0]) + " |"]
                        md_lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")
                        for row in rows[1:]:
                            md_lines.append("| " + " | ".join(row) + " |")
                        texts.append("\n".join(md_lines))
            if texts:
                content = "\n".join(texts)
                tokens = tokenizer.encode(content)
                results.append(
                    TextChunkSchema(
                        tokens=len(tokens),
                        content=content,
                        chunk_order_index=slide_num - 1,
                        full_doc_id="",
                        section_path=f"slide-{slide_num}",
                        page_number=slide_num,
                        modality="text",
                    )
                )
    except ImportError:
        logger.warning("python-pptx not available")
    return results


async def parse_xlsx(path: str) -> list[TextChunkSchema]:
    results: list[TextChunkSchema] = []
    tokenizer = TiktokenTokenizer()
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            header = [str(c) if c is not None else "" for c in rows[0]]
            md_lines = ["| " + " | ".join(header) + " |"]
            md_lines.append("| " + " | ".join("---" for _ in header) + " |")
            for row in rows[1:]:
                md_lines.append("| " + " | ".join(str(c) if c is not None else "" for c in row) + " |")
            md_text = "\n".join(md_lines)
            tokens = tokenizer.encode(md_text)
            results.append(
                TextChunkSchema(
                    tokens=len(tokens),
                    content=md_text,
                    chunk_order_index=sheet_idx,
                    full_doc_id="",
                    section_path=f"sheet-{sheet_name}",
                    page_number=None,
                    modality="table",
                )
            )
        wb.close()
    except ImportError:
        logger.warning("openpyxl not available")
    return results


# ═══════════════════════════════════════════════════════════════════════════
# SAFETY PRE-CHECKS
# ═══════════════════════════════════════════════════════════════════════════


def pii_scan(
    chunks: list[TextChunkSchema], policy: str = "flag"
) -> tuple[list[TextChunkSchema], list[dict]]:
    all_findings: list[dict] = []
    processed: list[TextChunkSchema] = []

    for chunk in chunks:
        findings = detect_pii(chunk["content"], policy=policy)
        if findings:
            all_findings.extend(
                [{**f, "chunk_order_index": chunk["chunk_order_index"]} for f in findings]
            )

        if policy == "reject" and findings:
            continue
        elif policy == "mask" and findings:
            content = chunk["content"]
            for f in sorted(findings, key=lambda x: x["start"], reverse=True):
                content = content[: f["start"]] + "[REDACTED]" + content[f["end"]:]
            chunk = {**chunk, "content": content}
        elif policy == "redact" and findings:
            content = chunk["content"]
            for f in sorted(findings, key=lambda x: x["start"], reverse=True):
                content = content[: f["start"]] + content[f["end"]:]
            chunk = {**chunk, "content": content.strip()}

        processed.append(chunk)

    return processed, all_findings


def prompt_injection_scan(
    chunks: list[TextChunkSchema], action: str = "flag"
) -> tuple[list[TextChunkSchema], list[dict]]:
    all_flags: list[dict] = []
    processed: list[TextChunkSchema] = []

    for chunk in chunks:
        result = detect_prompt_injection(chunk["content"])
        if result["patterns_matched"]:
            all_flags.append(
                {**result, "chunk_order_index": chunk["chunk_order_index"]}
            )

        if action == "reject" and result["risk_score"] > 0.5:
            continue

        processed.append(chunk)

    return processed, all_flags


# ═══════════════════════════════════════════════════════════════════════════
# ENTITY EXTRACTION
# ═══════════════════════════════════════════════════════════════════════════


def _truncate_entity_identifier(name: str) -> str:
    return name.strip().title()[:256]


def _parse_json_safe(text: str) -> Any:
    if json_repair is not None:
        return json_repair.loads(text)
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        match = re.search(r'[\[{]', text)
        if match:
            try:
                return json.loads(text[match.start():])
            except json.JSONDecodeError:
                return {}
        return {}


def _process_extraction_result(
    result: str | ExtractionResult,
    chunk_id: str,
    confidence_threshold: float = 0.3,
) -> tuple[dict, list]:
    nodes: dict[str, dict] = {}
    edges: list[dict] = []

    if isinstance(result, ExtractionResult):
        for ent in result.entities:
            conf = max(0.0, min(1.0, ent.confidence))
            if conf < confidence_threshold:
                continue
            name = _truncate_entity_identifier(ent.name)
            if not name:
                continue
            nodes[name] = {
                "type": ent.type,
                "description": ent.description,
                "source_id": chunk_id,
                "confidence": conf,
            }
        for rel in result.relations:
            conf = max(0.0, min(1.0, rel.confidence))
            if conf < confidence_threshold:
                continue
            src = _truncate_entity_identifier(rel.source)
            tgt = _truncate_entity_identifier(rel.target)
            if not src or not tgt:
                continue
            edges.append({
                "src_id": src,
                "tgt_id": tgt,
                "keywords": rel.keywords,
                "description": rel.description,
                "confidence": conf,
                "source_id": chunk_id,
                "valid_from": rel.valid_from,
                "valid_to": rel.valid_to,
            })
        return nodes, edges

    # Delimiter-based string parsing
    if not isinstance(result, str):
        return nodes, edges

    records = result.split(DEFAULT_RECORD_DELIMITER)
    for record in records:
        record = record.strip()
        if not record or DEFAULT_COMPLETION_DELIMITER in record:
            continue
        record = record.strip("()")
        fields = [f.strip() for f in record.split(DEFAULT_TUPLE_DELIMITER)]
        if len(fields) < 2:
            continue

        record_type = fields[0].lower().strip()
        if record_type == "entity" and len(fields) >= 5:
            name = _truncate_entity_identifier(fields[1])
            etype = fields[2].strip()
            desc = fields[3].strip()
            try:
                conf = max(0.0, min(1.0, float(fields[4])))
            except (ValueError, IndexError):
                conf = 0.5
            if conf < confidence_threshold or not name:
                continue
            nodes[name] = {
                "type": etype,
                "description": desc,
                "source_id": chunk_id,
                "confidence": conf,
            }
        elif record_type == "relationship" and len(fields) >= 6:
            src = _truncate_entity_identifier(fields[1])
            tgt = _truncate_entity_identifier(fields[2])
            keywords = fields[3].strip()
            desc = fields[4].strip()
            try:
                conf = max(0.0, min(1.0, float(fields[5])))
            except (ValueError, IndexError):
                conf = 0.5
            valid_from = fields[6].strip() if len(fields) > 6 and fields[6].strip() else None
            valid_to = fields[7].strip() if len(fields) > 7 and fields[7].strip() else None
            if conf < confidence_threshold or not src or not tgt:
                continue
            edges.append({
                "src_id": src,
                "tgt_id": tgt,
                "keywords": keywords,
                "description": desc,
                "confidence": conf,
                "source_id": chunk_id,
                "valid_from": valid_from,
                "valid_to": valid_to,
            })

    return nodes, edges


async def _extract_from_chunk(
    chunk: TextChunkSchema,
    global_config: dict,
    llm_response_cache: BaseKVStorage | None = None,
) -> tuple[dict, list]:
    llm_func = global_config["llm_model_func"]
    addon = global_config.get("addon_params", {})
    language = addon.get("language", "English")
    entity_types = addon.get("entity_types", "Use defaults")
    if isinstance(entity_types, list):
        entity_types = ", ".join(entity_types)
    confidence_threshold = global_config.get("confidence_threshold", 0.3)
    chunk_id = compute_mdhash_id(chunk["content"])
    use_structured = global_config.get("use_structured_output", True)

    if use_structured:
        system_prompt = PROMPTS["entity_extraction_structured_instructions"]
        user_prompt = PROMPTS["entity_extraction_user_prompt"].format(
            entity_types=entity_types, input_text=chunk["content"]
        )
        try:
            raw_result = await llm_func(
                user_prompt,
                system_prompt=system_prompt,
                response_schema=ExtractionResult,
            )
            if isinstance(raw_result, ExtractionResult):
                nodes, edges = _process_extraction_result(raw_result, chunk_id, confidence_threshold)
            elif isinstance(raw_result, str):
                parsed = _parse_json_safe(raw_result)
                if isinstance(parsed, dict):
                    extraction = ExtractionResult(
                        entities=[ExtractedEntity(**e) for e in parsed.get("entities", [])],
                        relations=[ExtractedRelation(**r) for r in parsed.get("relations", [])],
                    )
                    nodes, edges = _process_extraction_result(extraction, chunk_id, confidence_threshold)
                else:
                    nodes, edges = {}, []
            else:
                nodes, edges = {}, []
        except Exception:
            logger.debug("Structured extraction failed, falling back to delimiter-based")
            use_structured = False

    if not use_structured:
        system_prompt = PROMPTS["entity_extraction_system_prompt"].format(
            entity_types=entity_types,
            language=language,
            tuple_delimiter=DEFAULT_TUPLE_DELIMITER,
            record_delimiter=DEFAULT_RECORD_DELIMITER,
            completion_delimiter=DEFAULT_COMPLETION_DELIMITER,
            examples=PROMPTS.get("entity_extraction_examples", ""),
        )
        user_prompt = PROMPTS["entity_extraction_user_prompt"].format(
            entity_types=entity_types, input_text=chunk["content"]
        )
        if llm_response_cache and global_config.get("enable_llm_cache", True):
            raw_text, _ = await use_llm_func_with_cache(
                user_prompt, llm_func, system_prompt=system_prompt,
                llm_response_cache=llm_response_cache, cache_type="extraction",
            )
        else:
            raw_text = await llm_func(user_prompt, system_prompt=system_prompt)
        nodes, edges = _process_extraction_result(raw_text, chunk_id, confidence_threshold)

    # Gleaning passes
    max_gleaning = global_config.get("entity_extract_max_gleaning", 0)
    for _ in range(max_gleaning):
        continue_prompt = PROMPTS["entity_continue_extraction_user_prompt"]
        if llm_response_cache and global_config.get("enable_llm_cache", True):
            glean_text, _ = await use_llm_func_with_cache(
                continue_prompt, llm_func, system_prompt=system_prompt,
                llm_response_cache=llm_response_cache, cache_type="extraction_glean",
            )
        else:
            glean_text = await llm_func(continue_prompt, system_prompt=system_prompt)
        glean_nodes, glean_edges = _process_extraction_result(glean_text, chunk_id, confidence_threshold)
        for k, v in glean_nodes.items():
            if k not in nodes:
                nodes[k] = v
        edges.extend(glean_edges)
        if DEFAULT_COMPLETION_DELIMITER in glean_text:
            break

    return nodes, edges


async def extract_entities(
    chunks: list[TextChunkSchema],
    global_config: dict,
    pipeline_status: dict | None = None,
    llm_response_cache: BaseKVStorage | None = None,
) -> tuple[dict, list]:
    all_nodes: dict[str, dict] = {}
    all_edges: list[dict] = []

    tasks = [
        _extract_from_chunk(chunk, global_config, llm_response_cache)
        for chunk in chunks
    ]
    results = await asyncio.gather(*tasks, return_exceptions=True)

    for i, result in enumerate(results):
        if isinstance(result, Exception):
            logger.error(f"Entity extraction failed for chunk {i}: {result}")
            if pipeline_status is not None:
                pipeline_status.setdefault("errors", []).append(str(result))
            continue
        chunk_nodes, chunk_edges = result
        for name, data in chunk_nodes.items():
            if name in all_nodes:
                existing = all_nodes[name]
                existing["description"] = existing["description"] + " | " + data["description"]
                existing["source_id"] = existing["source_id"] + "," + data["source_id"]
                existing["confidence"] = max(existing["confidence"], data["confidence"])
            else:
                all_nodes[name] = data
        all_edges.extend(chunk_edges)
        if pipeline_status is not None:
            pipeline_status["chunks_processed"] = pipeline_status.get("chunks_processed", 0) + 1

    return all_nodes, all_edges


# ═══════════════════════════════════════════════════════════════════════════
# ENTITY RESOLUTION
# ═══════════════════════════════════════════════════════════════════════════


async def resolve_entities(
    candidate_entities: dict,
    knowledge_graph: BaseGraphStorage,
    entities_vdb: BaseVectorStorage,
    global_config: dict,
) -> list[list[str]]:
    threshold = global_config.get("resolution_similarity_threshold", 0.85)
    llm_func = global_config["llm_model_func"]
    merge_groups: list[list[str]] = []
    entity_names = list(candidate_entities.keys())
    visited_pairs: set[tuple[str, str]] = set()

    try:
        from rapidfuzz import fuzz
    except ImportError:
        fuzz = None  # type: ignore[assignment]

    for name in entity_names:
        try:
            similar = await entities_vdb.query(name, top_k=5)
        except Exception:
            continue

        for match in similar:
            match_name = match.get("entity_name", match.get("id", ""))
            if not match_name or match_name == name:
                continue
            pair = tuple(sorted([name, match_name]))
            if pair in visited_pairs:
                continue
            visited_pairs.add(pair)

            sim_score = match.get("similarity", match.get("distance", 0.0))
            if sim_score < threshold:
                continue

            # High similarity + high edit similarity → auto-merge
            if fuzz is not None:
                edit_ratio = fuzz.ratio(name.lower(), match_name.lower()) / 100.0
                if edit_ratio > 0.9 and sim_score > threshold:
                    _add_to_merge_group(merge_groups, name, match_name)
                    continue

            # Ambiguous: call LLM adjudicator
            existing_node = await knowledge_graph.get_node(match_name)
            candidate_data = candidate_entities.get(name, {})
            prompt = PROMPTS["entity_resolution_adjudicator"].format(
                entity_a_name=name,
                entity_a_type=candidate_data.get("type", ""),
                entity_a_description=candidate_data.get("description", ""),
                entity_a_context=candidate_data.get("source_id", ""),
                entity_b_name=match_name,
                entity_b_type=(existing_node or {}).get("type", ""),
                entity_b_description=(existing_node or {}).get("description", ""),
                entity_b_context=(existing_node or {}).get("source_id", ""),
            )
            try:
                response = await llm_func(prompt)
                if "SAME" in response.upper().split("\n")[0]:
                    _add_to_merge_group(merge_groups, name, match_name)
            except Exception as e:
                logger.debug(f"Entity resolution LLM call failed: {e}")

    return merge_groups


def _add_to_merge_group(groups: list[list[str]], a: str, b: str) -> None:
    for group in groups:
        if a in group or b in group:
            if a not in group:
                group.append(a)
            if b not in group:
                group.append(b)
            return
    groups.append([a, b])


# ═══════════════════════════════════════════════════════════════════════════
# MERGE & UPSERT
# ═══════════════════════════════════════════════════════════════════════════


async def _handle_entity_relation_summary(
    description_list: list[str],
    description_type: str,
    description_name: str,
    global_config: dict,
    llm_response_cache: BaseKVStorage | None = None,
) -> str:
    if len(description_list) <= 1:
        return description_list[0] if description_list else ""

    language = global_config.get("addon_params", {}).get("language", "English")
    llm_func = global_config["llm_model_func"]
    prompt = PROMPTS["summarize_entity_descriptions"].format(
        description_type=description_type,
        description_name=description_name,
        description_list="\n".join(f"- {d}" for d in description_list),
        summary_length="a concise paragraph",
        language=language,
    )
    if llm_response_cache and global_config.get("enable_llm_cache", True):
        summary, _ = await use_llm_func_with_cache(
            prompt, llm_func, llm_response_cache=llm_response_cache,
            cache_type="summary",
        )
    else:
        summary = await llm_func(prompt)
    return summary


async def _merge_nodes_then_upsert(
    nodes: dict,
    knowledge_graph: BaseGraphStorage,
    entities_vdb: BaseVectorStorage,
    lexical_storage: BaseLexicalStorage,
    global_config: dict,
    llm_response_cache: BaseKVStorage | None = None,
) -> None:
    embedding_func = global_config["embedding_func"]

    for name, data in nodes.items():
        existing = await knowledge_graph.get_node(name)
        if existing:
            old_descs = existing.get("description", "").split(" | ")
            new_descs = data["description"].split(" | ")
            all_descs = [d.strip() for d in old_descs + new_descs if d.strip()]
            unique_descs = list(dict.fromkeys(all_descs))

            if len(unique_descs) > 1:
                merged_desc = await _handle_entity_relation_summary(
                    unique_descs, "entity", name, global_config, llm_response_cache
                )
            else:
                merged_desc = unique_descs[0] if unique_descs else ""

            old_sources = set(existing.get("source_id", "").split(","))
            new_sources = set(data.get("source_id", "").split(","))
            merged_sources = ",".join(old_sources | new_sources)

            node_data = {
                "type": data.get("type") or existing.get("type", ""),
                "description": merged_desc,
                "source_id": merged_sources,
                "confidence": max(
                    existing.get("confidence", 0.0), data.get("confidence", 0.0)
                ),
            }
        else:
            node_data = {
                "type": data.get("type", ""),
                "description": data.get("description", ""),
                "source_id": data.get("source_id", ""),
                "confidence": data.get("confidence", 0.5),
            }

        await knowledge_graph.upsert_node(name, node_data)

        try:
            embeddings = await embedding_func([f"{name}: {node_data['description']}"])
            embed_vec = embeddings[0] if len(embeddings) > 0 else None
            await entities_vdb.upsert({
                compute_mdhash_id(name): {
                    "entity_name": name,
                    "content": f"{name}: {node_data['description']}",
                    "embedding": embed_vec,
                }
            })
        except Exception as e:
            logger.warning(f"VDB upsert failed for entity {name}: {e}")

        try:
            await lexical_storage.upsert({
                compute_mdhash_id(name): {
                    "entity_name": name,
                    "content": f"{name} ({node_data['type']}): {node_data['description']}",
                }
            })
        except Exception as e:
            logger.debug(f"Lexical upsert failed for entity {name}: {e}")


async def _merge_edges_then_upsert(
    edges: list,
    knowledge_graph: BaseGraphStorage,
    relationships_vdb: BaseVectorStorage,
    global_config: dict,
    llm_response_cache: BaseKVStorage | None = None,
) -> None:
    embedding_func = global_config["embedding_func"]

    for edge in edges:
        src = edge["src_id"]
        tgt = edge["tgt_id"]
        existing = await knowledge_graph.get_edge(src, tgt)

        if existing:
            old_descs = existing.get("description", "").split(" | ")
            new_descs = edge["description"].split(" | ")
            all_descs = [d.strip() for d in old_descs + new_descs if d.strip()]
            unique_descs = list(dict.fromkeys(all_descs))

            if len(unique_descs) > 1:
                merged_desc = await _handle_entity_relation_summary(
                    unique_descs, "relationship", f"{src} -> {tgt}",
                    global_config, llm_response_cache,
                )
            else:
                merged_desc = unique_descs[0] if unique_descs else ""

            old_sources = set(existing.get("source_id", "").split(","))
            new_sources = set(edge.get("source_id", "").split(","))

            edge_data = {
                "keywords": edge.get("keywords", existing.get("keywords", "")),
                "description": merged_desc,
                "source_id": ",".join(old_sources | new_sources),
                "confidence": max(
                    existing.get("confidence", 0.0), edge.get("confidence", 0.0)
                ),
                "valid_from": edge.get("valid_from") or existing.get("valid_from"),
                "valid_to": edge.get("valid_to") or existing.get("valid_to"),
            }
        else:
            edge_data = {
                "keywords": edge.get("keywords", ""),
                "description": edge.get("description", ""),
                "source_id": edge.get("source_id", ""),
                "confidence": edge.get("confidence", 0.5),
                "valid_from": edge.get("valid_from"),
                "valid_to": edge.get("valid_to"),
            }

        await knowledge_graph.upsert_edge(src, tgt, edge_data)

        try:
            content = f"{src} -> {tgt}: {edge_data['description']}"
            embeddings = await embedding_func([content])
            embed_vec = embeddings[0] if len(embeddings) > 0 else None
            edge_id = compute_mdhash_id(f"{src}-{tgt}")
            await relationships_vdb.upsert({
                edge_id: {
                    "src_id": src,
                    "tgt_id": tgt,
                    "content": content,
                    "embedding": embed_vec,
                }
            })
        except Exception as e:
            logger.warning(f"VDB upsert failed for edge {src}->{tgt}: {e}")


# ═══════════════════════════════════════════════════════════════════════════
# COMMUNITY DETECTION
# ═══════════════════════════════════════════════════════════════════════════


async def build_communities(
    knowledge_graph: BaseGraphStorage,
    global_config: dict,
    llm_response_cache: BaseKVStorage | None = None,
) -> dict:
    levels = global_config.get("community_levels", 3)
    llm_func = global_config["llm_model_func"]

    community_hierarchy = await knowledge_graph.detect_communities("leiden", levels=levels)

    communities = community_hierarchy.get("communities", {})
    for community_id, community_data in communities.items():
        members = community_data.get("members", [])
        if not members:
            continue

        entities_text = []
        relations_text = []
        for member in members:
            node = await knowledge_graph.get_node(member)
            if node:
                entities_text.append(
                    f"- {member} ({node.get('type', 'Unknown')}): {node.get('description', '')}"
                )
            node_edges = await knowledge_graph.get_node_edges(member)
            for src, tgt in node_edges:
                edge = await knowledge_graph.get_edge(src, tgt)
                if edge:
                    relations_text.append(
                        f"- {src} -> {tgt}: {edge.get('description', '')}"
                    )

        prompt = PROMPTS["community_report"].format(
            entities="\n".join(entities_text) or "None",
            relations="\n".join(relations_text) or "None",
        )

        if llm_response_cache and global_config.get("enable_llm_cache", True):
            report_text, _ = await use_llm_func_with_cache(
                prompt, llm_func, llm_response_cache=llm_response_cache,
                cache_type="community_report",
            )
        else:
            report_text = await llm_func(prompt)

        parsed = _parse_json_safe(report_text)
        if isinstance(parsed, dict):
            community_data["report"] = parsed
        else:
            community_data["report"] = {"title": f"Community {community_id}", "summary": report_text}

    return community_hierarchy


# ═══════════════════════════════════════════════════════════════════════════
# QUERY FUNCTIONS – Helpers
# ═══════════════════════════════════════════════════════════════════════════


async def rewrite_query(
    query: str, conversation_history: list[dict], global_config: dict
) -> str:
    if not conversation_history:
        return query
    llm_func = global_config["llm_model_func"]
    history_text = "\n".join(
        f"{msg.get('role', 'user')}: {msg.get('content', '')}"
        for msg in conversation_history
    )
    prompt = PROMPTS["query_rewrite"].format(
        conversation_history=history_text, query=query
    )
    result = await llm_func(prompt)
    return result.strip() or query


async def maybe_decompose(query: str, global_config: dict) -> list[str]:
    llm_func = global_config["llm_model_func"]
    prompt = PROMPTS["query_decomposition"].format(query=query)
    result = await llm_func(prompt)
    parsed = _parse_json_safe(result)
    if isinstance(parsed, list) and len(parsed) > 0:
        return [str(q) for q in parsed]
    return [query]


async def maybe_hyde(query: str, global_config: dict) -> str:
    llm_func = global_config["llm_model_func"]
    prompt = PROMPTS["hyde_generation"].format(query=query)
    result = await llm_func(prompt)
    return result.strip()


async def get_keywords_from_query(
    query: str,
    global_config: dict,
    llm_response_cache: BaseKVStorage | None = None,
) -> dict:
    llm_func = global_config["llm_model_func"]
    prompt = PROMPTS["keywords_extraction"].format(query=query)
    examples = PROMPTS.get("keywords_extraction_examples", "")
    full_prompt = prompt + "\n\n" + examples if examples else prompt

    if llm_response_cache and global_config.get("enable_llm_cache", True):
        result, _ = await use_llm_func_with_cache(
            full_prompt, llm_func, llm_response_cache=llm_response_cache,
            cache_type="keywords",
        )
    else:
        result = await llm_func(full_prompt)

    parsed = _parse_json_safe(result)
    if isinstance(parsed, dict):
        return {
            "high_level_keywords": parsed.get("high_level_keywords", []),
            "low_level_keywords": parsed.get("low_level_keywords", []),
        }
    return {"high_level_keywords": [], "low_level_keywords": []}


def _apply_token_truncation(context: dict, param: QueryParam) -> dict:
    if context.get("entities"):
        context["entities"] = truncate_list_by_token_size(
            context["entities"], key="content",
            max_token_size=param.max_entity_tokens,
        )
    if context.get("relations"):
        context["relations"] = truncate_list_by_token_size(
            context["relations"], key="content",
            max_token_size=param.max_relation_tokens,
        )
    # Enforce total token budget
    all_items = []
    for key in ("entities", "relations", "chunks", "communities"):
        for item in context.get(key, []):
            all_items.append({**item, "_source_key": key})
    truncated = truncate_list_by_token_size(
        all_items, key="content", max_token_size=param.max_total_tokens,
    )
    rebuilt: dict[str, list] = {"entities": [], "relations": [], "chunks": [], "communities": []}
    for item in truncated:
        src_key = item.pop("_source_key", "chunks")
        rebuilt.setdefault(src_key, []).append(item)
    context.update(rebuilt)
    return context


def apply_acl_filter(
    results: list[dict], user_id: str | None, user_groups: list[str]
) -> list[dict]:
    if user_id is None:
        return results
    filtered = []
    for item in results:
        acl = item.get("acl_policy")
        if acl is None:
            filtered.append(item)
            continue
        if isinstance(acl, ACLPolicy):
            if acl.can_access(user_id, user_groups):
                filtered.append(item)
        elif isinstance(acl, dict):
            policy = ACLPolicy(**acl)
            if policy.can_access(user_id, user_groups):
                filtered.append(item)
        else:
            filtered.append(item)
    return filtered


def apply_temporal_filter(results: list[dict], snapshot_at) -> list[dict]:
    if snapshot_at is None:
        return results
    if isinstance(snapshot_at, str):
        snapshot_at = datetime.fromisoformat(snapshot_at)
    filtered = []
    for item in results:
        valid_from = item.get("valid_from")
        valid_to = item.get("valid_to")
        if valid_from is not None:
            if isinstance(valid_from, str):
                valid_from = datetime.fromisoformat(valid_from)
            if snapshot_at < valid_from:
                continue
        if valid_to is not None:
            if isinstance(valid_to, str):
                valid_to = datetime.fromisoformat(valid_to)
            if snapshot_at > valid_to:
                continue
        filtered.append(item)
    return filtered


async def _run_grounded_check(
    answer: str, contexts: list[dict], global_config: dict
) -> list[dict]:
    verifier_func = global_config.get("verifier_func") or global_config.get("llm_model_func")
    if not verifier_func:
        return []

    claims = [s.strip() for s in re.split(r'(?<=[.!?])\s+', answer) if s.strip()]
    if not claims:
        return []

    context_text = "\n\n".join(
        c.get("content", str(c)) for c in contexts
    )

    results: list[dict] = []
    for claim in claims:
        prompt = PROMPTS["grounded_check"].format(
            claim=claim, context=context_text
        )
        try:
            response = await verifier_func(prompt)
            parsed = _parse_json_safe(response)
            if isinstance(parsed, dict):
                results.append({
                    "claim": claim,
                    "score": float(parsed.get("score", 0.0)),
                    "supporting_span": parsed.get("supporting_span", ""),
                })
            else:
                results.append({"claim": claim, "score": 0.0, "supporting_span": ""})
        except Exception:
            results.append({"claim": claim, "score": 0.0, "supporting_span": ""})

    return results


# ═══════════════════════════════════════════════════════════════════════════
# QUERY CONTEXT BUILDING
# ═══════════════════════════════════════════════════════════════════════════


async def _build_query_context(
    query: str,
    keywords: dict,
    param: QueryParam,
    global_config: dict,
    knowledge_graph: BaseGraphStorage,
    entities_vdb: BaseVectorStorage,
    relationships_vdb: BaseVectorStorage,
    chunks_vdb: BaseVectorStorage,
    lexical_storage: BaseLexicalStorage,
) -> dict:
    context: dict[str, list] = {
        "entities": [],
        "relations": [],
        "chunks": [],
        "communities": [],
        "references": [],
    }
    rerank_func = global_config.get("rerank_func")
    ref_counter = [0]

    def _make_ref(source: str) -> str:
        ref_counter[0] += 1
        return f"ref-{ref_counter[0]}"

    acl_filter = None
    if param.user_id:
        acl_filter = {"user_id": param.user_id, "user_groups": param.user_groups}

    async def _search_entities(q: str, top_k: int) -> list[dict]:
        results = await entities_vdb.query(q, top_k=top_k, acl_filter=acl_filter)
        return results

    async def _search_relations(q: str, top_k: int) -> list[dict]:
        results = await relationships_vdb.query(q, top_k=top_k, acl_filter=acl_filter)
        return results

    async def _search_chunks(q: str, top_k: int) -> list[dict]:
        results = await chunks_vdb.query(q, top_k=top_k, acl_filter=acl_filter)
        return results

    async def _bm25_search(q: str, top_k: int) -> list[dict]:
        try:
            return await lexical_storage.search_bm25(q, top_k=top_k)
        except Exception:
            return []

    all_keywords = keywords.get("low_level_keywords", []) + keywords.get("high_level_keywords", [])
    keyword_query = " ".join(all_keywords) if all_keywords else query

    if param.mode == "local":
        entity_results = await _search_entities(query, param.top_k)
        if param.enable_hybrid_lexical:
            bm25_entities = await _bm25_search(keyword_query, param.top_k)
            entity_results = reciprocal_rank_fusion(
                [entity_results, bm25_entities], k=param.rrf_k
            )
        for ent in entity_results:
            ent_name = ent.get("entity_name", ent.get("id", ""))
            ref_id = _make_ref("entity")
            ent["ref_id"] = ref_id
            ent.setdefault("content", ent_name)
            context["entities"].append(ent)
            context["references"].append({"id": ref_id, "source": f"entity: {ent_name}"})
            # Expand to connected relations
            if ent_name:
                edges = await knowledge_graph.get_node_edges(ent_name)
                for src, tgt in edges[:10]:
                    edge_data = await knowledge_graph.get_edge(src, tgt)
                    if edge_data:
                        r_ref = _make_ref("relation")
                        rel_item = {
                            "src_id": src, "tgt_id": tgt,
                            "content": f"{src} -> {tgt}: {edge_data.get('description', '')}",
                            "ref_id": r_ref, **edge_data,
                        }
                        context["relations"].append(rel_item)
                        context["references"].append({"id": r_ref, "source": f"relation: {src} -> {tgt}"})
        # Gather source chunks
        source_ids: set[str] = set()
        for ent in context["entities"]:
            for sid in ent.get("source_id", "").split(","):
                if sid.strip():
                    source_ids.add(sid.strip())
        if source_ids:
            chunk_results = await _search_chunks(query, param.chunk_top_k)
            for ch in chunk_results:
                c_ref = _make_ref("chunk")
                ch["ref_id"] = c_ref
                ch.setdefault("content", "")
                context["chunks"].append(ch)
                context["references"].append({"id": c_ref, "source": f"chunk: {ch.get('id', '')}"})

    elif param.mode == "global":
        rel_results = await _search_relations(query, param.top_k)
        if param.enable_hybrid_lexical:
            bm25_rels = await _bm25_search(keyword_query, param.top_k)
            rel_results = reciprocal_rank_fusion(
                [rel_results, bm25_rels], k=param.rrf_k
            )
        connected_entities: set[str] = set()
        for rel in rel_results:
            r_ref = _make_ref("relation")
            rel["ref_id"] = r_ref
            rel.setdefault("content", rel.get("description", ""))
            context["relations"].append(rel)
            context["references"].append({"id": r_ref, "source": f"relation: {rel.get('src_id', '')} -> {rel.get('tgt_id', '')}"})
            connected_entities.add(rel.get("src_id", ""))
            connected_entities.add(rel.get("tgt_id", ""))
        for ent_name in connected_entities:
            if not ent_name:
                continue
            node = await knowledge_graph.get_node(ent_name)
            if node:
                e_ref = _make_ref("entity")
                ent_item = {
                    "entity_name": ent_name,
                    "content": f"{ent_name}: {node.get('description', '')}",
                    "ref_id": e_ref, **node,
                }
                context["entities"].append(ent_item)
                context["references"].append({"id": e_ref, "source": f"entity: {ent_name}"})
        chunk_results = await _search_chunks(query, param.chunk_top_k)
        for ch in chunk_results:
            c_ref = _make_ref("chunk")
            ch["ref_id"] = c_ref
            ch.setdefault("content", "")
            context["chunks"].append(ch)
            context["references"].append({"id": c_ref, "source": f"chunk: {ch.get('id', '')}"})

    elif param.mode == "hybrid":
        local_task = _build_query_context(
            query, keywords,
            QueryParam(**{**param.__dict__, "mode": "local"}),
            global_config, knowledge_graph, entities_vdb,
            relationships_vdb, chunks_vdb, lexical_storage,
        )
        global_task = _build_query_context(
            query, keywords,
            QueryParam(**{**param.__dict__, "mode": "global"}),
            global_config, knowledge_graph, entities_vdb,
            relationships_vdb, chunks_vdb, lexical_storage,
        )
        local_ctx, global_ctx = await asyncio.gather(local_task, global_task)
        for key in ("entities", "relations", "chunks", "communities", "references"):
            seen_ids: set[str] = set()
            for item in local_ctx.get(key, []) + global_ctx.get(key, []):
                item_id = item.get("id", item.get("ref_id", id(item)))
                if item_id not in seen_ids:
                    seen_ids.add(item_id)
                    context[key].append(item)

    elif param.mode == "naive":
        chunk_results = await _search_chunks(query, param.top_k)
        if param.enable_hybrid_lexical:
            bm25_chunks = await _bm25_search(query, param.top_k)
            chunk_results = reciprocal_rank_fusion(
                [chunk_results, bm25_chunks], k=param.rrf_k
            )
        for ch in chunk_results:
            c_ref = _make_ref("chunk")
            ch["ref_id"] = c_ref
            ch.setdefault("content", "")
            context["chunks"].append(ch)
            context["references"].append({"id": c_ref, "source": f"chunk: {ch.get('id', '')}"})

    elif param.mode == "mix":
        kg_task = _build_query_context(
            query, keywords,
            QueryParam(**{**param.__dict__, "mode": "local"}),
            global_config, knowledge_graph, entities_vdb,
            relationships_vdb, chunks_vdb, lexical_storage,
        )
        naive_task = _build_query_context(
            query, keywords,
            QueryParam(**{**param.__dict__, "mode": "naive"}),
            global_config, knowledge_graph, entities_vdb,
            relationships_vdb, chunks_vdb, lexical_storage,
        )
        kg_ctx, naive_ctx = await asyncio.gather(kg_task, naive_task)
        for key in ("entities", "relations", "chunks", "communities", "references"):
            seen_ids: set[str] = set()
            for item in kg_ctx.get(key, []) + naive_ctx.get(key, []):
                item_id = item.get("id", item.get("ref_id", id(item)))
                if item_id not in seen_ids:
                    seen_ids.add(item_id)
                    context[key].append(item)

    # Rerank if available
    if param.enable_rerank and rerank_func and context.get("chunks"):
        try:
            reranked = await rerank_func(
                query, [c.get("content", "") for c in context["chunks"]]
            )
            if isinstance(reranked, list):
                reranked_chunks = []
                for item in reranked:
                    idx = item.get("index", 0) if isinstance(item, dict) else item
                    if isinstance(idx, int) and idx < len(context["chunks"]):
                        reranked_chunks.append(context["chunks"][idx])
                if reranked_chunks:
                    context["chunks"] = reranked_chunks
        except Exception as e:
            logger.debug(f"Reranking failed: {e}")

    # Apply temporal filter
    if param.snapshot_at:
        context["relations"] = apply_temporal_filter(context["relations"], param.snapshot_at)

    # Apply ACL filter
    if param.user_id:
        for key in ("entities", "relations", "chunks"):
            context[key] = apply_acl_filter(context[key], param.user_id, param.user_groups)

    context = _apply_token_truncation(context, param)
    return context


# ═══════════════════════════════════════════════════════════════════════════
# QUERY FUNCTIONS – Main Entry Points
# ═══════════════════════════════════════════════════════════════════════════


async def kg_query(
    query: str,
    param: QueryParam,
    global_config: dict,
    knowledge_graph: BaseGraphStorage,
    entities_vdb: BaseVectorStorage,
    relationships_vdb: BaseVectorStorage,
    chunks_vdb: BaseVectorStorage,
    lexical_storage: BaseLexicalStorage,
    llm_response_cache: BaseKVStorage | None = None,
) -> QueryResult:
    import time as _time
    t0 = _time.monotonic()
    llm_func = param.model_func or global_config["llm_model_func"]

    effective_query = query
    if param.conversation_history:
        effective_query = await rewrite_query(query, param.conversation_history, global_config)

    sub_queries = await maybe_decompose(effective_query, global_config)

    aggregated_context: dict[str, list] = {
        "entities": [], "relations": [], "chunks": [], "communities": [], "references": [],
    }

    for sq in sub_queries:
        kw = await get_keywords_from_query(sq, global_config, llm_response_cache)
        ctx = await _build_query_context(
            sq, kw, param, global_config, knowledge_graph,
            entities_vdb, relationships_vdb, chunks_vdb, lexical_storage,
        )
        for key in aggregated_context:
            aggregated_context[key].extend(ctx.get(key, []))

    # Deduplicate references
    seen_refs: set[str] = set()
    unique_refs = []
    for ref in aggregated_context["references"]:
        if ref["id"] not in seen_refs:
            seen_refs.add(ref["id"])
            unique_refs.append(ref)
    aggregated_context["references"] = unique_refs

    # Format context
    context_data = PROMPTS["kg_query_context"].format(
        entities="\n".join(
            f"[{e.get('ref_id', '')}] {e.get('content', '')}" for e in aggregated_context["entities"]
        ) or "None",
        relations="\n".join(
            f"[{r.get('ref_id', '')}] {r.get('content', '')}" for r in aggregated_context["relations"]
        ) or "None",
        chunks="\n".join(
            f"[{c.get('ref_id', '')}] {c.get('content', '')}" for c in aggregated_context["chunks"]
        ) or "None",
        communities="\n".join(
            f"[{cm.get('ref_id', '')}] {cm.get('content', '')}" for cm in aggregated_context["communities"]
        ) or "None",
        references="\n".join(
            f"{r['id']}: {r.get('source', '')}" for r in aggregated_context["references"]
        ) or "None",
    )

    prompt = PROMPTS["rag_response"].format(
        response_type=param.response_type,
        user_prompt=effective_query,
        context_data=context_data,
    )

    if param.only_need_context:
        return QueryResult(
            content="", raw_data=aggregated_context,
            references=aggregated_context["references"],
            latency_ms=(_time.monotonic() - t0) * 1000,
        )

    if param.only_need_prompt:
        return QueryResult(
            content=prompt, raw_data=aggregated_context,
            latency_ms=(_time.monotonic() - t0) * 1000,
        )

    if param.stream:
        async def _stream():
            async for chunk in llm_func(prompt, stream=True):
                yield chunk
        return QueryResult(
            response_iterator=_stream(), is_streaming=True,
            references=aggregated_context["references"],
            communities_used=[c.get("id", "") for c in aggregated_context["communities"]],
            latency_ms=(_time.monotonic() - t0) * 1000,
        )

    answer = await llm_func(prompt)
    answer = sanitize_output(answer)

    grounded = []
    if param.verifier_enabled and (global_config.get("verifier_func") or global_config.get("llm_model_func")):
        all_contexts = aggregated_context["entities"] + aggregated_context["relations"] + aggregated_context["chunks"]
        if all_contexts:
            grounded = await _run_grounded_check(answer, all_contexts, global_config)

    return QueryResult(
        content=answer,
        raw_data=aggregated_context,
        references=aggregated_context["references"],
        grounded_check=grounded,
        communities_used=[c.get("id", "") for c in aggregated_context["communities"]],
        latency_ms=(_time.monotonic() - t0) * 1000,
    )


async def naive_query(
    query: str,
    param: QueryParam,
    global_config: dict,
    chunks_vdb: BaseVectorStorage,
    lexical_storage: BaseLexicalStorage,
    llm_response_cache: BaseKVStorage | None = None,
) -> QueryResult:
    import time as _time
    t0 = _time.monotonic()
    llm_func = param.model_func or global_config["llm_model_func"]

    effective_query = query
    if param.conversation_history:
        effective_query = await rewrite_query(query, param.conversation_history, global_config)

    # HyDE for short queries
    search_query = effective_query
    tokenizer = TiktokenTokenizer()
    if len(tokenizer.encode(effective_query)) < 10:
        try:
            search_query = await maybe_hyde(effective_query, global_config)
        except Exception:
            search_query = effective_query

    acl_filter = None
    if param.user_id:
        acl_filter = {"user_id": param.user_id, "user_groups": param.user_groups}

    # Hybrid vector + BM25
    vector_results = await chunks_vdb.query(search_query, top_k=param.top_k, acl_filter=acl_filter)
    bm25_results = []
    if param.enable_hybrid_lexical:
        try:
            bm25_results = await lexical_storage.search_bm25(effective_query, top_k=param.top_k)
        except Exception:
            pass

    if bm25_results:
        fused = reciprocal_rank_fusion([vector_results, bm25_results], k=param.rrf_k)
    else:
        fused = vector_results

    # Rerank
    rerank_func = global_config.get("rerank_func")
    if param.enable_rerank and rerank_func and fused:
        try:
            reranked = await rerank_func(
                effective_query, [c.get("content", "") for c in fused]
            )
            if isinstance(reranked, list):
                reranked_chunks = []
                for item in reranked:
                    idx = item.get("index", 0) if isinstance(item, dict) else item
                    if isinstance(idx, int) and idx < len(fused):
                        reranked_chunks.append(fused[idx])
                if reranked_chunks:
                    fused = reranked_chunks
        except Exception as e:
            logger.debug(f"Reranking failed in naive_query: {e}")

    # ACL + temporal filter
    if param.user_id:
        fused = apply_acl_filter(fused, param.user_id, param.user_groups)
    if param.snapshot_at:
        fused = apply_temporal_filter(fused, param.snapshot_at)

    # Truncate
    fused = truncate_list_by_token_size(
        [{"content": c.get("content", ""), **c} for c in fused],
        key="content", max_token_size=param.max_total_tokens,
    )

    # Build references
    references = []
    for i, ch in enumerate(fused):
        ref_id = f"ref-{i + 1}"
        ch["ref_id"] = ref_id
        references.append({"id": ref_id, "source": f"chunk: {ch.get('id', '')}"})

    context_data = PROMPTS["naive_query_context"].format(
        chunks="\n".join(
            f"[{c.get('ref_id', '')}] {c.get('content', '')}" for c in fused
        ) or "None",
        references="\n".join(
            f"{r['id']}: {r.get('source', '')}" for r in references
        ) or "None",
    )

    prompt = PROMPTS["naive_rag_response"].format(
        response_type=param.response_type,
        user_prompt=effective_query,
        context_data=context_data,
    )

    if param.only_need_context:
        return QueryResult(
            content="", raw_data={"chunks": fused},
            references=references,
            latency_ms=(_time.monotonic() - t0) * 1000,
        )

    if param.stream:
        async def _stream():
            async for chunk in llm_func(prompt, stream=True):
                yield chunk
        return QueryResult(
            response_iterator=_stream(), is_streaming=True,
            references=references,
            latency_ms=(_time.monotonic() - t0) * 1000,
        )

    answer = await llm_func(prompt)
    answer = sanitize_output(answer)

    grounded = []
    if param.verifier_enabled and fused:
        grounded = await _run_grounded_check(answer, fused, global_config)

    return QueryResult(
        content=answer,
        raw_data={"chunks": fused},
        references=references,
        grounded_check=grounded,
        latency_ms=(_time.monotonic() - t0) * 1000,
    )


async def community_query(
    query: str,
    param: QueryParam,
    global_config: dict,
    knowledge_graph: BaseGraphStorage,
    llm_response_cache: BaseKVStorage | None = None,
) -> QueryResult:
    import time as _time
    t0 = _time.monotonic()
    llm_func = param.model_func or global_config["llm_model_func"]

    keywords = await get_keywords_from_query(query, global_config, llm_response_cache)
    all_kw = keywords.get("high_level_keywords", []) + keywords.get("low_level_keywords", [])
    kw_text = " ".join(all_kw) if all_kw else query

    # Get community hierarchy
    community_hierarchy = await knowledge_graph.detect_communities(
        "leiden", levels=global_config.get("community_levels", 3)
    )
    communities = community_hierarchy.get("communities", {})

    # Score communities by keyword overlap
    scored: list[tuple[str, float, dict]] = []
    for cid, cdata in communities.items():
        report = cdata.get("report", {})
        summary = report.get("summary", "") if isinstance(report, dict) else str(report)
        title = report.get("title", "") if isinstance(report, dict) else ""
        combined = f"{title} {summary}".lower()
        score = sum(1 for kw in all_kw if kw.lower() in combined)
        if score > 0 or not all_kw:
            scored.append((cid, score, cdata))

    scored.sort(key=lambda x: x[1], reverse=True)
    top_communities = scored[:param.top_k]

    # Assemble context
    community_texts = []
    references = []
    communities_used = []
    for i, (cid, _, cdata) in enumerate(top_communities):
        report = cdata.get("report", {})
        if isinstance(report, dict):
            text = f"Community: {report.get('title', cid)}\n{report.get('summary', '')}"
            findings = report.get("findings", [])
            for f in findings:
                text += f"\n- {f.get('explanation', '')}"
        else:
            text = str(report)
        ref_id = f"ref-{i + 1}"
        community_texts.append(f"[{ref_id}] {text}")
        references.append({"id": ref_id, "source": f"community: {cid}"})
        communities_used.append(cid)

    context_data = "\n\n".join(community_texts) or "No relevant communities found."

    prompt = PROMPTS["rag_response"].format(
        response_type=param.response_type,
        user_prompt=query,
        context_data=context_data,
    )

    if param.only_need_context:
        return QueryResult(
            content="", raw_data={"communities": [c[2] for c in top_communities]},
            references=references, communities_used=communities_used,
            latency_ms=(_time.monotonic() - t0) * 1000,
        )

    answer = await llm_func(prompt)
    answer = sanitize_output(answer)

    return QueryResult(
        content=answer,
        raw_data={"communities": [c[2] for c in top_communities]},
        references=references,
        communities_used=communities_used,
        latency_ms=(_time.monotonic() - t0) * 1000,
    )
